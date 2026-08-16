"""
MCP Office 模块 — Word / Excel / PowerPoint 文档操作
注册为 AI 可调用的工具，支持创建、编辑和保存 Office 文档
"""

import json
import os
import re
import tempfile
import time
from pathlib import Path

from flask import jsonify, request, send_file

from . import tool_registry

# ============================================================
# 会话管理 — 每个用户在服务端维护一个工作目录
# 会话在服务运行期间保持有效，重启后清空
# ============================================================
_workspace = Path(tempfile.gettempdir()) / "mcp_office_workspace"
_workspace.mkdir(parents=True, exist_ok=True)

_sessions = {}  # session_id -> { word_file, excel_file, ppt_file, ... }


def _get_or_create_session(session_id: str) -> dict:
    """获取或创建会话。同一 session_id 的 Word/Excel/PPT 共享同一个会话。"""
    if session_id not in _sessions:
        sub = _workspace / session_id
        sub.mkdir(parents=True, exist_ok=True)
        _sessions[session_id] = {
            "dir": str(sub),
            "word_doc": None,       # python-docx Document 对象
            "word_path": None,
            "excel_wb": None,       # openpyxl Workbook 对象
            "excel_path": None,
            "ppt_prs": None,        # python-pptx Presentation 对象
            "ppt_path": None,
        }
    return _sessions[session_id]


# ============================================================
# Word 工具定义
# ============================================================
WORD_CREATE_DEF = {
    "name": "word_create",
    "description": (
        "创建一个新的 Word 文档。"
        "⚠️ 仅在用户明确要求创建/生成/导出/起草 Word 文档或报告时使用。"
        "不要在用户只是闲聊或询问信息时自动创建文档。"
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "session_id": {"type": "string", "description": "会话 ID，用于跟踪同一用户的文档"},
        },
        "required": [],
    },
}

WORD_ADD_HEADING_DEF = {
    "name": "word_add_heading",
    "description": "向当前 Word 文档添加标题。",
    "parameters": {
        "type": "object",
        "properties": {
            "text": {"type": "string", "description": "标题文本"},
            "level": {"type": "integer", "description": "标题级别 1-9，1 为最高级", "default": 1},
            "session_id": {"type": "string", "description": "会话 ID"},
        },
        "required": ["text"],
    },
}

WORD_ADD_PARAGRAPH_DEF = {
    "name": "word_add_paragraph",
    "description": "向当前 Word 文档添加一个段落。可以设置加粗、斜体等格式。",
    "parameters": {
        "type": "object",
        "properties": {
            "text": {"type": "string", "description": "段落文本"},
            "bold": {"type": "boolean", "description": "是否加粗", "default": False},
            "italic": {"type": "boolean", "description": "是否斜体", "default": False},
            "font_size": {"type": "integer", "description": "字体大小（磅）", "default": 12},
            "session_id": {"type": "string", "description": "会话 ID"},
        },
        "required": ["text"],
    },
}

WORD_ADD_TABLE_DEF = {
    "name": "word_add_table",
    "description": "向当前 Word 文档添加一个表格。rows 是二维数组，每行是一个数组。",
    "parameters": {
        "type": "object",
        "properties": {
            "rows": {"type": "array", "description": "表格数据，二维数组 [[列1,列2,...], ...]"},
            "session_id": {"type": "string", "description": "会话 ID"},
        },
        "required": ["rows"],
    },
}

WORD_SAVE_DEF = {
    "name": "word_save",
    "description": "保存当前 Word 文档到指定路径。返回保存后的文件路径。",
    "parameters": {
        "type": "object",
        "properties": {
            "filename": {"type": "string", "description": "文件名（不含路径），如 'report.docx'"},
            "session_id": {"type": "string", "description": "会话 ID"},
        },
        "required": ["filename"],
    },
}

# ============================================================
# Excel 工具定义
# ============================================================
EXCEL_CREATE_DEF = {
    "name": "excel_create",
    "description": (
        "创建一个新的 Excel 工作簿。"
        "⚠️ 仅在用户明确要求创建/生成/导出 Excel 文件、表格文件、电子表格时才使用。"
        "不要在用户只是说'列出'、'列表'、'整理'时自动创建 Excel —— 用普通文本回复即可。"
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "session_id": {"type": "string", "description": "会话 ID"},
        },
        "required": [],
    },
}

EXCEL_WRITE_CELL_DEF = {
    "name": "excel_write_cell",
    "description": "向 Excel 工作表写入单元格内容。可指定工作表名、单元格位置和值。",
    "parameters": {
        "type": "object",
        "properties": {
            "sheet": {"type": "string", "description": "工作表名称", "default": "Sheet"},
            "cell": {"type": "string", "description": "单元格位置，如 'A1', 'B3'"},
            "value": {"type": "string", "description": "要写入的值（字符串或数字）"},
            "session_id": {"type": "string", "description": "会话 ID"},
        },
        "required": ["cell", "value"],
    },
}

EXCEL_READ_CELL_DEF = {
    "name": "excel_read_cell",
    "description": "读取 Excel 中指定单元格的值。",
    "parameters": {
        "type": "object",
        "properties": {
            "sheet": {"type": "string", "description": "工作表名称", "default": "Sheet"},
            "cell": {"type": "string", "description": "单元格位置，如 'A1'"},
            "session_id": {"type": "string", "description": "会话 ID"},
        },
        "required": ["cell"],
    },
}

EXCEL_ADD_SHEET_DEF = {
    "name": "excel_add_sheet",
    "description": "向当前 Excel 工作簿添加一个新的工作表。",
    "parameters": {
        "type": "object",
        "properties": {
            "name": {"type": "string", "description": "新工作表名称"},
            "session_id": {"type": "string", "description": "会话 ID"},
        },
        "required": ["name"],
    },
}

EXCEL_SAVE_DEF = {
    "name": "excel_save",
    "description": "保存当前 Excel 工作簿到指定路径。返回保存后的文件路径。",
    "parameters": {
        "type": "object",
        "properties": {
            "filename": {"type": "string", "description": "文件名（不含路径），如 'data.xlsx'"},
            "session_id": {"type": "string", "description": "会话 ID"},
        },
        "required": ["filename"],
    },
}

# ============================================================
# PowerPoint 工具定义
# ============================================================
PPT_CREATE_DEF = {
    "name": "ppt_create",
    "description": (
        "创建一个新的 PowerPoint 演示文稿。"
        "⚠️ 仅在用户明确要求创建/生成/制作 PPT、演示文稿、幻灯片时使用。"
        "不要在用户只是闲聊或询问信息时自动创建演示文稿。"
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "session_id": {"type": "string", "description": "会话 ID"},
        },
        "required": [],
    },
}

PPT_ADD_SLIDE_DEF = {
    "name": "ppt_add_slide",
    "description": "向演示文稿添加一页幻灯片。可选择布局类型。",
    "parameters": {
        "type": "object",
        "properties": {
            "title": {"type": "string", "description": "幻灯片标题"},
            "layout": {"type": "string", "description": "布局类型: title（标题页）, content（标题+内容）, blank（空白）", "default": "content"},
            "session_id": {"type": "string", "description": "会话 ID"},
        },
        "required": ["title"],
    },
}

PPT_ADD_TEXT_DEF = {
    "name": "ppt_add_text",
    "description": "向当前幻灯片添加文本框。",
    "parameters": {
        "type": "object",
        "properties": {
            "text": {"type": "string", "description": "文本内容"},
            "left": {"type": "number", "description": "左边距，厘米", "default": 2},
            "top": {"type": "number", "description": "上边距，厘米", "default": 5},
            "width": {"type": "number", "description": "宽度，厘米", "default": 20},
            "height": {"type": "number", "description": "高度，厘米", "default": 5},
            "font_size": {"type": "integer", "description": "字体大小（磅）", "default": 18},
            "session_id": {"type": "string", "description": "会话 ID"},
        },
        "required": ["text"],
    },
}

PPT_ADD_BULLET_LIST_DEF = {
    "name": "ppt_add_bullet_list",
    "description": "向当前幻灯片添加一个项目符号列表。",
    "parameters": {
        "type": "object",
        "properties": {
            "items": {"type": "array", "description": "列表项，字符串数组"},
            "left": {"type": "number", "description": "左边距，厘米", "default": 2},
            "top": {"type": "number", "description": "上边距，厘米", "default": 5},
            "session_id": {"type": "string", "description": "会话 ID"},
        },
        "required": ["items"],
    },
}

PPT_SAVE_DEF = {
    "name": "ppt_save",
    "description": "保存当前 PowerPoint 演示文稿到指定路径。",
    "parameters": {
        "type": "object",
        "properties": {
            "filename": {"type": "string", "description": "文件名（不含路径），如 'presentation.pptx'"},
            "session_id": {"type": "string", "description": "会话 ID"},
        },
        "required": ["filename"],
    },
}


# ============================================================
# 辅助函数 — 统一 session + 文档获取，消除重复的检查逻辑
# ============================================================
def _sid(args: dict) -> str:
    return args.get("session_id", "default")


def _get_doc(args: dict) -> tuple:
    """返回 (session, doc) 或 (session, error_str)"""
    session = _get_or_create_session(_sid(args))
    doc = session.get("word_doc")
    if doc is None:
        return session, "❌ 还未创建 Word 文档，请先调用 word_create。"
    return session, doc


def _get_wb(args: dict) -> tuple:
    session = _get_or_create_session(_sid(args))
    wb = session.get("excel_wb")
    if wb is None:
        return session, "❌ 还未创建 Excel 工作簿，请先调用 excel_create。"
    return session, wb


def _get_prs(args: dict) -> tuple:
    session = _get_or_create_session(_sid(args))
    prs = session.get("ppt_prs")
    if prs is None:
        return session, "❌ 还未创建 PPT，请先调用 ppt_create。"
    return session, prs


def _get_slide(args: dict) -> tuple:
    session, result = _get_prs(args)
    if isinstance(result, str):
        return session, result
    slide = session.get("ppt_current_slide")
    if slide is None:
        return session, "❌ 还没有当前幻灯片，请先调用 ppt_add_slide。"
    return session, slide


# ============================================================
# Word 执行器
# ============================================================
def _exec_word_create(args: dict) -> str:
    from docx import Document
    session = _get_or_create_session(_sid(args))
    doc = Document()
    session["word_doc"] = doc
    session["word_path"] = None
    return "✅ Word 文档已创建，可以开始编辑。请使用 word_add_heading、word_add_paragraph 等工具添加内容。"


def _exec_word_add_heading(args: dict) -> str:
    session, result = _get_doc(args)
    if isinstance(result, str): return result
    text = args.get("text", "")
    level = min(max(int(args.get("level", 1)), 1), 9)
    result.add_heading(text, level=level)
    return f"✅ 已添加 {level} 级标题: {text}"


def _exec_word_add_paragraph(args: dict) -> str:
    session, result = _get_doc(args)
    if isinstance(result, str): return result
    text = args.get("text", "")
    p = result.add_paragraph()
    run = p.add_run(text)
    run.bold = bool(args.get("bold", False))
    run.italic = bool(args.get("italic", False))
    run.font.size = int(args.get("font_size", 12)) * 12700
    return f"✅ 已添加段落: {text[:80]}"


def _exec_word_add_table(args: dict) -> str:
    session, result = _get_doc(args)
    if isinstance(result, str): return result
    rows = args.get("rows", [])
    if not rows:
        return "❌ 表格数据不能为空"
    table = result.add_table(rows=len(rows), cols=len(rows[0]))
    table.style = "Light Grid Accent 1"
    for i, row_data in enumerate(rows):
        for j, cell_text in enumerate(row_data):
            table.cell(i, j).text = str(cell_text)
    return f"✅ 已添加 {len(rows)}x{len(rows[0])} 表格"


def _exec_word_save(args: dict) -> str:
    session, result = _get_doc(args)
    if isinstance(result, str): return result
    filename = args.get("filename", "document.docx")
    filepath = Path(session["dir"]) / filename
    result.save(str(filepath))
    session["word_path"] = str(filepath)
    return f"✅ Word 文档已保存到: {filepath}"


# ============================================================
# Excel 执行器
# ============================================================
def _exec_excel_create(args: dict) -> str:
    from openpyxl import Workbook
    session = _get_or_create_session(_sid(args))
    wb = Workbook()
    session["excel_wb"] = wb
    session["excel_path"] = None
    return "✅ Excel 工作簿已创建，默认工作表名为 'Sheet'。可以使用 excel_write_cell 写入数据。"


def _exec_excel_write_cell(args: dict) -> str:
    session, result = _get_wb(args)
    if isinstance(result, str): return result
    sheet_name = args.get("sheet", "Sheet")
    if sheet_name not in result.sheetnames:
        return f"❌ 工作表 '{sheet_name}' 不存在。可用工作表: {', '.join(result.sheetnames)}"
    cell, value = args.get("cell", ""), args.get("value", "")
    result[sheet_name][cell] = value
    return f"✅ 已将 {cell} 设为: {value}"


def _exec_excel_read_cell(args: dict) -> str:
    session, result = _get_wb(args)
    if isinstance(result, str): return result
    sheet_name = args.get("sheet", "Sheet")
    if sheet_name not in result.sheetnames:
        return f"❌ 工作表 '{sheet_name}' 不存在。"
    cell = args.get("cell", "")
    value = result[sheet_name][cell].value
    return f"📊 {cell} = {value}"


def _exec_excel_add_sheet(args: dict) -> str:
    session, result = _get_wb(args)
    if isinstance(result, str): return result
    name = args.get("name", "")
    result.create_sheet(title=name)
    return f"✅ 已添加工作表: {name}"


def _exec_excel_save(args: dict) -> str:
    session, result = _get_wb(args)
    if isinstance(result, str): return result
    filename = args.get("filename", "workbook.xlsx")
    filepath = Path(session["dir"]) / filename
    result.save(str(filepath))
    session["excel_path"] = str(filepath)
    return f"✅ Excel 工作簿已保存到: {filepath}"


# ============================================================
# PowerPoint 执行器
# ============================================================
def _exec_ppt_create(args: dict) -> str:
    from pptx import Presentation
    session = _get_or_create_session(_sid(args))
    prs = Presentation()
    prs.slide_width = 9144000 * 2
    prs.slide_height = 5143500 * 2
    session["ppt_prs"] = prs
    session["ppt_path"] = None
    session["ppt_current_slide"] = None
    return "✅ PowerPoint 演示文稿已创建。可以使用 ppt_add_slide 添加幻灯片。"


def _exec_ppt_add_slide(args: dict) -> str:
    from pptx import Presentation
    session, result = _get_prs(args)
    if isinstance(result, str): return result
    title = args.get("title", "")
    layout_map = {"title": 0, "content": 1, "blank": 6}
    layout_idx = layout_map.get(args.get("layout", "content"), 1)
    slide = result.slides.add_slide(result.slide_layouts[layout_idx])
    if slide.shapes.title:
        slide.shapes.title.text = title
    session["ppt_current_slide"] = slide
    return f"✅ 已添加幻灯片: {title}"


def _exec_ppt_add_text(args: dict) -> str:
    from pptx.util import Inches, Pt
    session, result = _get_slide(args)
    if isinstance(result, str): return result
    text = args.get("text", "")
    txBox = result.shapes.add_textbox(
        Inches(float(args.get("left", 2))), Inches(float(args.get("top", 5))),
        Inches(float(args.get("width", 20))), Inches(float(args.get("height", 5))),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(int(args.get("font_size", 18)))
    return f"✅ 已添加文本框: {text[:80]}"


def _exec_ppt_add_bullet_list(args: dict) -> str:
    from pptx.util import Inches, Pt
    session, result = _get_slide(args)
    if isinstance(result, str): return result
    items = args.get("items", [])
    txBox = result.shapes.add_textbox(
        Inches(float(args.get("left", 2))), Inches(float(args.get("top", 5))),
        Inches(20), Inches(10),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
    return f"✅ 已添加项目符号列表，共 {len(items)} 项"


def _exec_ppt_save(args: dict) -> str:
    session, result = _get_prs(args)
    if isinstance(result, str): return result
    filename = args.get("filename", "presentation.pptx")
    filepath = Path(session["dir"]) / filename
    result.save(str(filepath))
    session["ppt_path"] = str(filepath)
    return f"✅ PPT 演示文稿已保存到: {filepath}"


# ============================================================
# 注册所有工具到 tool_registry
# ============================================================
_tool_list = [
    # Word
    (WORD_CREATE_DEF, _exec_word_create),
    (WORD_ADD_HEADING_DEF, _exec_word_add_heading),
    (WORD_ADD_PARAGRAPH_DEF, _exec_word_add_paragraph),
    (WORD_ADD_TABLE_DEF, _exec_word_add_table),
    (WORD_SAVE_DEF, _exec_word_save),
    # Excel
    (EXCEL_CREATE_DEF, _exec_excel_create),
    (EXCEL_WRITE_CELL_DEF, _exec_excel_write_cell),
    (EXCEL_READ_CELL_DEF, _exec_excel_read_cell),
    (EXCEL_ADD_SHEET_DEF, _exec_excel_add_sheet),
    (EXCEL_SAVE_DEF, _exec_excel_save),
    # PowerPoint
    (PPT_CREATE_DEF, _exec_ppt_create),
    (PPT_ADD_SLIDE_DEF, _exec_ppt_add_slide),
    (PPT_ADD_TEXT_DEF, _exec_ppt_add_text),
    (PPT_ADD_BULLET_LIST_DEF, _exec_ppt_add_bullet_list),
    (PPT_SAVE_DEF, _exec_ppt_save),
]

for tool_def, executor in _tool_list:
    tool_registry.register(tool_def["name"], tool_def, executor)


# ============================================================
# Flask 路由
# ============================================================
def register_routes(app, http_session=None):
    @app.route("/api/mcp/office/download/<session_id>/<filename>")
    def download_file(session_id, filename):
        """下载生成的 Office 文件"""
        session = _get_or_create_session(session_id)
        filepath = Path(session["dir"]) / filename
        if not filepath.exists():
            return jsonify({"error": "文件不存在或已被清理"}), 404
        return send_file(str(filepath), as_attachment=True, download_name=filename)

    @app.route("/api/mcp/office/list-files/<session_id>")
    def list_files(session_id):
        """列出会话中生成的文件"""
        session = _get_or_create_session(session_id)
        files = []
        for f in Path(session["dir"]).iterdir():
            if f.is_file():
                files.append({
                    "name": f.name,
                    "size": f.stat().st_size,
                    "modified": time.strftime("%Y-%m-%d %H:%M:%S", time.localtime(f.stat().st_mtime)),
                })
        return jsonify({"files": files, "dir": session["dir"]})

    @app.route("/api/mcp/office/delete/<session_id>/<filename>", methods=["DELETE"])
    def delete_file(session_id, filename):
        """删除工作区中的文件（保存到本地后清理服务端副本）"""
        session = _get_or_create_session(session_id)
        filepath = Path(session["dir"]) / filename
        if filepath.exists():
            try:
                filepath.unlink()
            except OSError as e:
                return jsonify({"error": f"删除失败: {e}"}), 500
        # 清除对应的路径引用
        ext = Path(filename).suffix
        if ext == ".docx" and session.get("word_path") == str(filepath):
            session["word_path"] = None
        elif ext == ".xlsx" and session.get("excel_path") == str(filepath):
            session["excel_path"] = None
        elif ext == ".pptx" and session.get("ppt_path") == str(filepath):
            session["ppt_path"] = None
        return jsonify({"ok": True})

    @app.route("/api/mcp/office/workspace-info/<session_id>")
    def workspace_info(session_id):
        """获取工作区状态"""
        session = _get_or_create_session(session_id)
        return jsonify({
            "word_active": session.get("word_doc") is not None,
            "excel_active": session.get("excel_wb") is not None,
            "ppt_active": session.get("ppt_prs") is not None,
            "word_path": session.get("word_path"),
            "excel_path": session.get("excel_path"),
            "ppt_path": session.get("ppt_path"),
            "workspace_dir": session["dir"],
        })
